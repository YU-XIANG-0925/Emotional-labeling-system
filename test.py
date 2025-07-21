import os
from pptx import Presentation

def export_notes_to_txt(pptx_filename: str, output_filename: str):
    """
    開啟指定的 PowerPoint (.pptx) 檔案，讀取每一頁投影片的備忘錄，
    並將其依照指定格式匯出至 .txt 檔案。

    Args:
        pptx_filename (str): 來源 .pptx 檔案的名稱。
        output_filename (str): 欲儲存的 .txt 檔案的名稱。
    """
    # 檢查來源檔案是否存在
    if not os.path.exists(pptx_filename):
        print(f"錯誤：找不到檔案 '{pptx_filename}'，請確認檔案名稱與路徑是否正確。")
        return

    try:
        # 載入簡報檔案
        prs = Presentation(pptx_filename)
        
        # 準備一個列表來存放所有要寫入的文字行
        output_lines = []

        print(f"開始讀取 '{pptx_filename}'...")

        # 使用 enumerate 來同時取得索引 (頁碼) 和投影片物件
        # 設定 start=1 讓頁碼從 1 開始，符合一般人的習慣
        for i, slide in enumerate(prs.slides, 1):
            
            # 建立符合您要求的分頁標頭
            page_header = f"===第{i}頁==="
            output_lines.append(page_header)
            
            notes_text = ""
            # 步驟 1: 檢查該投影片是否有備忘錄頁 (notes_slide)
            if slide.has_notes_slide:
                # 步驟 2: 取得備忘錄的文字框 (notes_text_frame) 並讀取文字
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            
            # 步驟 3: 判斷讀取到的文字是否為空
            if notes_text:
                # 如果有文字，則加入列表
                output_lines.append(notes_text)
            else:
                # 如果沒有文字或沒有備忘錄頁，則使用指定的預設文字
                output_lines.append("(無對應講稿)")
        
        # 使用換行符號將所有內容串接成一個完整的字串
        final_content = "\n".join(output_lines)
        
        # 將最終內容寫入指定的 .txt 檔案，使用 utf-8 編碼以支援中文
        with open(output_filename, 'w', encoding='utf-8') as f:
            f.write(final_content)
            
        print(f"成功！所有備忘錄已儲存至 '{output_filename}'。")
        print(f"總共處理了 {len(prs.slides)} 頁投影片。")

    except Exception as e:
        print(f"處理過程中發生錯誤：{e}")


if __name__ == '__main__':
    # --- 請在此處設定您的檔案名稱 ---
    # 您的 PowerPoint 檔案名稱
    PPTX_FILE = r"D:\OneDrive - 逢甲大學\Documents\星期一\250519_GestureCoach.pptx"
    # 您希望輸出的文字檔案名稱
    OUTPUT_TXT_FILE = r"D:\OneDrive - 逢甲大學\Documents\星期一\250519_GestureCoach.txt"
    
    export_notes_to_txt(PPTX_FILE, OUTPUT_TXT_FILE)